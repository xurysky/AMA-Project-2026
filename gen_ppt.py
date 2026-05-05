#!/usr/bin/env python3
"""Generate AMA Capstone Project Final Presentation PPT (25 pages)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color palette ──────────────────────────────────────────────
DEEP_BLUE   = RGBColor(0x00, 0x20, 0x50)   # #002050
GOLD        = RGBColor(0xFF, 0xD7, 0x00)   # #FFD700
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF0, 0xF2, 0xF5)
DARK_TEXT   = RGBColor(0x33, 0x33, 0x33)
MID_BLUE    = RGBColor(0x00, 0x3D, 0x7A)
LIGHT_BLUE  = RGBColor(0xE8, 0xEE, 0xF6)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

# ── Helper functions ───────────────────────────────────────────

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=16,
                          color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                          line_spacing=1.3, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * 0.4)
    return txBox

def add_bullet_textbox(slide, left, top, width, height, bullets, font_size=16,
                       color=DARK_TEXT, bullet_char='•', font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet_char} {b}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * 0.5)
    return txBox

def add_title_bar(slide, title, subtitle=None):
    """Add a dark blue title bar at top with optional gold subtitle"""
    bar = add_rect(slide, Inches(0), Inches(0), W, Inches(1.3), DEEP_BLUE)
    # Gold accent line
    add_rect(slide, Inches(0), Inches(1.3), W, Inches(0.06), GOLD)
    # Title text
    add_textbox(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.7),
                title, font_size=30, color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.75), Inches(11), Inches(0.4),
                    subtitle, font_size=14, color=GOLD, bold=False)

def add_page_number(slide, num):
    add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
                str(num), font_size=10, color=RGBColor(0x99,0x99,0x99),
                alignment=PP_ALIGN.RIGHT)

def content_slide(title, bullets, page_num, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    add_title_bar(slide, title, subtitle)
    add_bullet_textbox(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(5.2),
                       bullets, font_size=18, color=DARK_TEXT)
    add_page_number(slide, page_num)
    return slide

# ══════════════════════════════════════════════════════════════
# Page 1: Cover
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DEEP_BLUE)

# Gold accent line at top
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), GOLD)

# Project title
add_textbox(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(1.2),
            "Autonomous Retail Intelligence\nAgent Platform",
            font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide, Inches(1), Inches(3.3), Inches(11.3), Inches(0.8),
            "基于优衣库的多 Agent 智能零售平台",
            font_size=24, color=GOLD, bold=False, alignment=PP_ALIGN.CENTER)

# Divider
add_rect(slide, Inches(5), Inches(4.3), Inches(3.3), Inches(0.04), GOLD)

# Presenter info
add_multiline_textbox(slide, Inches(1), Inches(4.8), Inches(11.3), Inches(1.5),
    [
        "汇报人：圆子  |  Senior Solution Engineer",
        "AMA Capstone Project  |  Project 46"
    ], font_size=18, color=RGBColor(0xCC,0xCC,0xCC), alignment=PP_ALIGN.CENTER)

# Bottom gold line
add_rect(slide, Inches(0), Inches(7.42), W, Inches(0.08), GOLD)

# ══════════════════════════════════════════════════════════════
# Page 2: Table of Contents
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "目录")

toc_items = [
    ("01", "项目背景与挑战"),
    ("02", "解决方案概述"),
    ("03", "5 Agent 架构设计"),
    ("04", "技术架构"),
    ("05", "亚洲特殊性"),
    ("06", "业务场景演示"),
    ("07", "预期效果"),
    ("08", "实施路径"),
    ("09", "风险与应对"),
    ("10", "Q&A"),
]
for i, (num, item) in enumerate(toc_items):
    y = Inches(1.7 + i * 0.52)
    # Number
    add_textbox(slide, Inches(1.5), y, Inches(0.8), Inches(0.45),
                num, font_size=20, color=GOLD, bold=True)
    # Dot
    add_rect(slide, Inches(2.3), y + Inches(0.18), Inches(0.12), Inches(0.12), DEEP_BLUE)
    # Text
    add_textbox(slide, Inches(2.6), y, Inches(8), Inches(0.45),
                item, font_size=20, color=DARK_TEXT)
add_page_number(slide, 2)

# ══════════════════════════════════════════════════════════════
# Page 3: 项目背景
# ══════════════════════════════════════════════════════════════
content_slide("项目背景", [
    "全球零售业面临数字化转型挑战，传统模式难以为继",
    "优衣库（迅销集团）：全球 3,500+ 门店，年产量 13 亿件",
    "核心痛点：数据孤岛 → 各渠道数据无法打通",
    "核心痛点：个性化不足 → 千篇一律的客户体验",
    "核心痛点：库存浪费 → 行业平均库存浪费率超 30%",
    "核心痛点：定价滞后 → 无法快速响应市场变化",
], 3)

# ══════════════════════════════════════════════════════════════
# Page 4: 业务挑战
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "业务挑战")

challenges = [
    ("数据分散", "客户数据分散在 App、门店、电商、社交媒体等多个触点"),
    ("库存浪费", "行业平均库存浪费率 30%+，每年数十亿美元损失"),
    ("定价滞后", "定价策略无法实时响应市场变化和竞争动态"),
    ("营销低效", "营销活动缺乏精准触达，ROI 难以衡量"),
]
for i, (title, desc) in enumerate(challenges):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.0)
    box = add_rounded_rect(slide, x, y, Inches(2.8), Inches(4.5), LIGHT_BLUE, DEEP_BLUE)
    # Icon area (gold circle)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), y + Inches(0.4), Inches(1), Inches(1))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GOLD
    circle.line.fill.background()
    add_textbox(slide, x + Inches(0.9), y + Inches(0.55), Inches(1), Inches(0.7),
                str(i+1), font_size=28, color=DEEP_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.7), Inches(2.4), Inches(0.5),
                title, font_size=20, color=DEEP_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), y + Inches(2.3), Inches(2.4), Inches(1.8),
                desc, font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 4)

# ══════════════════════════════════════════════════════════════
# Page 5: 解决方案概述
# ══════════════════════════════════════════════════════════════
content_slide("解决方案概述", [
    "5 个 AI Agent 协同工作，覆盖零售全链路",
    "统一客户数据平台，打通 App / 门店 / 电商 / 社交媒体",
    "自主决策 + 人工审批，兼顾效率与风控",
    "Azure 全栈技术支撑，企业级安全与可扩展性",
    "亚洲市场特殊性深度适配（隐私、支付、电商生态）",
], 5)

# ══════════════════════════════════════════════════════════════
# Page 6: 5 Agent 架构图
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "5 Agent 架构设计")

agents = [
    ("Customer\nUnderstanding", "客户理解", Inches(0.5)),
    ("Personalization", "个性化", Inches(3.0)),
    ("Inventory", "库存管理", Inches(5.5)),
    ("Pricing", "动态定价", Inches(8.0)),
    ("Marketing", "营销优化", Inches(10.5)),
]
for name_cn, name_en, x in agents:
    y = Inches(2.2)
    box = add_rounded_rect(slide, x, y, Inches(2.2), Inches(1.8), DEEP_BLUE)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.3), Inches(2.0), Inches(1.2),
                name_cn, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), y + Inches(1.1), Inches(2.0), Inches(0.5),
                name_en, font_size=12, color=GOLD, alignment=PP_ALIGN.CENTER)

# Arrows between agents
for i in range(4):
    x_start = agents[i][2] + Inches(2.2)
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x_start, Inches(2.8), Inches(0.8), Inches(0.5))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = GOLD
    arrow.line.fill.background()

# Description below
add_multiline_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5), [
    "数据流：Customer Understanding → Personalization → Inventory / Pricing / Marketing",
    "协作模式：同步调用（REST API）+ 异步事件（Event Hubs）",
    "状态共享：Cosmos DB 统一状态管理",
], font_size=16, color=DARK_TEXT)
add_page_number(slide, 6)

# ══════════════════════════════════════════════════════════════
# Page 7: Customer Understanding Agent
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "Customer Understanding Agent", "构建 360° 客户视图")

# Left: Info
add_bullet_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(5.0), [
    "职责：构建 360° 客户视图，统一多渠道数据",
    "技术：合成 Persona + Agent Memory + 向量搜索",
    "输入：App 行为 / 门店 RFID / 电商订单 / 会员数据",
    "输出：统一客户画像 + 行为预测 + 流失预警",
], font_size=17, color=DARK_TEXT)

# Right: Key metrics box
box = add_rounded_rect(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(4.5), LIGHT_BLUE, DEEP_BLUE)
add_textbox(slide, Inches(7.3), Inches(1.9), Inches(5.0), Inches(0.5),
            "核心能力", font_size=20, color=DEEP_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_textbox(slide, Inches(7.3), Inches(2.6), Inches(5.0), Inches(3.5), [
    "合成 Persona 解决冷启动问题",
    "Agent Memory 实现跨会话记忆",
    "向量搜索实现语义级客户匹配",
    "实时行为追踪与预测",
], font_size=15, color=DARK_TEXT)
add_page_number(slide, 7)

# ══════════════════════════════════════════════════════════════
# Page 8: Personalization Agent
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "Personalization Agent", "跨渠道千人千面推荐")

add_bullet_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(5.0), [
    "职责：跨渠道千人千面个性化推荐",
    "技术：U-A-L 三支柱 + Hybrid AI",
    "输入：客户画像 + 实时行为 + 库存数据",
    "输出：个性化推荐 + 下一步最佳行动",
], font_size=17, color=DARK_TEXT)

box = add_rounded_rect(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(4.5), LIGHT_BLUE, DEEP_BLUE)
add_textbox(slide, Inches(7.3), Inches(1.9), Inches(5.0), Inches(0.5),
            "U-A-L 三支柱", font_size=20, color=DEEP_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_textbox(slide, Inches(7.3), Inches(2.6), Inches(5.0), Inches(3.5), [
    "Understanding：深度理解客户意图",
    "Action：智能推荐下一步行动",
    "Learning：持续学习优化推荐效果",
    "Hybrid AI：规则引擎 + 深度学习融合",
], font_size=15, color=DARK_TEXT)
add_page_number(slide, 8)

# ══════════════════════════════════════════════════════════════
# Page 9: Inventory Agent
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "Inventory Agent", "需求预测 + 自动补货")

add_bullet_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(5.0), [
    "职责：需求预测 + 自动补货决策",
    "技术：Prophet + LSTM 混合预测模型",
    "输入：历史销售 / 天气 / 促销 / 竞争数据",
    "输出：7 天需求预测 + 智能补货建议",
], font_size=17, color=DARK_TEXT)

box = add_rounded_rect(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(4.5), LIGHT_BLUE, DEEP_BLUE)
add_textbox(slide, Inches(7.3), Inches(1.9), Inches(5.0), Inches(0.5),
            "预测模型", font_size=20, color=DEEP_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_textbox(slide, Inches(7.3), Inches(2.6), Inches(5.0), Inches(3.5), [
    "Prophet：捕获季节性和趋势",
    "LSTM：学习复杂非线性模式",
    "混合模型：兼顾精度与可解释性",
    "实时调整：天气/事件触发重预测",
], font_size=15, color=DARK_TEXT)
add_page_number(slide, 9)

# ══════════════════════════════════════════════════════════════
# Page 10: Pricing Agent
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "Pricing Agent", "动态定价最大化利润")

add_bullet_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(5.0), [
    "职责：动态定价，最大化利润与周转率",
    "技术：Three-Layer Architecture（Data → ML → LLM）",
    "核心：价格弹性线性回归模型",
    "输出：最优价格 + 利润影响预测",
], font_size=17, color=DARK_TEXT)

box = add_rounded_rect(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(4.5), LIGHT_BLUE, DEEP_BLUE)
add_textbox(slide, Inches(7.3), Inches(1.9), Inches(5.0), Inches(0.5),
            "Three-Layer Architecture", font_size=20, color=DEEP_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_textbox(slide, Inches(7.3), Inches(2.6), Inches(5.0), Inches(3.5), [
    "Data Layer：清洗、聚合、特征工程",
    "ML Layer：价格弹性建模与优化",
    "LLM Layer：自然语言解释定价逻辑",
    "人工审批：关键价格变动需确认",
], font_size=15, color=DARK_TEXT)
add_page_number(slide, 10)

# ══════════════════════════════════════════════════════════════
# Page 11: Marketing Agent
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "Marketing Agent", "自主营销活动 + 预算优化")

add_bullet_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(5.0), [
    "职责：自主营销活动策划与预算优化",
    "技术：强化学习 + A/B 测试框架",
    "输入：客户画像 + 市场趋势 + 竞品动态",
    "输出：活动计划 + 预算分配 + 效果预测",
], font_size=17, color=DARK_TEXT)

box = add_rounded_rect(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(4.5), LIGHT_BLUE, DEEP_BLUE)
add_textbox(slide, Inches(7.3), Inches(1.9), Inches(5.0), Inches(0.5),
            "核心能力", font_size=20, color=DEEP_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_bullet_textbox(slide, Inches(7.3), Inches(2.6), Inches(5.0), Inches(3.5), [
    "强化学习：动态调整营销策略",
    "A/B 测试：数据驱动的效果验证",
    "预算优化：ROI 最大化分配",
    "多渠道协同：App/短信/社交媒体",
], font_size=15, color=DARK_TEXT)
add_page_number(slide, 11)

# ══════════════════════════════════════════════════════════════
# Page 12: Agent 间协作协议
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "Agent 间协作协议")

# Three boxes for collaboration mechanisms
mechanisms = [
    ("同步调用", "REST API 实时通信\n适用于需要即时响应的场景\n如：个性化推荐请求"),
    ("异步事件", "Event Hubs 事件驱动\n适用于非实时数据流\n如：库存更新、价格变动"),
    ("状态共享", "Cosmos DB 统一状态\nAgent 间共享上下文\n冲突解决：优先级 + 时间戳"),
]
for i, (title, desc) in enumerate(mechanisms):
    x = Inches(0.8 + i * 4.2)
    box = add_rounded_rect(slide, x, Inches(2.0), Inches(3.8), Inches(4.5), LIGHT_BLUE, DEEP_BLUE)
    add_textbox(slide, x + Inches(0.2), Inches(2.2), Inches(3.4), Inches(0.5),
                title, font_size=22, color=DEEP_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_rect(slide, x + Inches(0.5), Inches(2.8), Inches(2.8), Inches(0.04), GOLD)
    add_textbox(slide, x + Inches(0.3), Inches(3.1), Inches(3.2), Inches(3.0),
                desc, font_size=15, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 12)

# ══════════════════════════════════════════════════════════════
# Page 13: 技术架构（三层）
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "技术架构（三层）")

layers = [
    ("Application Layer", "LLM 推理层", "Azure AI Foundry + Agent Framework\nAzure OpenAI Service\n自然语言理解与生成", DEEP_BLUE),
    ("ML Layer", "预测模型层", "Azure ML + Prophet + LSTM\n价格弹性模型\n需求预测引擎", MID_BLUE),
    ("Data Layer", "数据基础层", "Microsoft Fabric + Cosmos DB\nEvent Hubs 实时流\n数据湖 + 特征存储", RGBColor(0x00, 0x55, 0x99)),
]
for i, (name, cn, desc, color) in enumerate(layers):
    y = Inches(1.7 + i * 1.85)
    box = add_rounded_rect(slide, Inches(0.8), y, Inches(11.5), Inches(1.6), color)
    add_textbox(slide, Inches(1.2), y + Inches(0.15), Inches(3.5), Inches(0.5),
                name, font_size=20, color=WHITE, bold=True)
    add_textbox(slide, Inches(1.2), y + Inches(0.6), Inches(3.5), Inches(0.4),
                cn, font_size=14, color=GOLD)
    add_textbox(slide, Inches(5.5), y + Inches(0.15), Inches(6.5), Inches(1.3),
                desc, font_size=14, color=WHITE)
add_page_number(slide, 13)

# ══════════════════════════════════════════════════════════════
# Page 14: Azure 技术栈
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "Azure 技术栈")

tech_items = [
    ("Azure AI Foundry", "Agent Framework 统一编排"),
    ("Azure OpenAI Service", "GPT-4o / GPT-4.1 推理引擎"),
    ("Microsoft Fabric", "统一数据分析平台"),
    ("Cosmos DB", "全球分布式状态存储"),
    ("Azure API Management", "MCP Gateway 统一网关"),
    ("Azure ML", "模型训练与部署"),
    ("Event Hubs", "实时事件流处理"),
]
for i, (name, desc) in enumerate(tech_items):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.3)
    y = Inches(1.7 + row * 1.35)
    box = add_rounded_rect(slide, x, y, Inches(5.8), Inches(1.15), LIGHT_BLUE, DEEP_BLUE)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.1), Inches(5.2), Inches(0.5),
                name, font_size=18, color=DEEP_BLUE, bold=True)
    add_textbox(slide, x + Inches(0.3), y + Inches(0.6), Inches(5.2), Inches(0.4),
                desc, font_size=14, color=DARK_TEXT)
add_page_number(slide, 14)

# ══════════════════════════════════════════════════════════════
# Page 15: 亚洲特殊性（1/2）
# ══════════════════════════════════════════════════════════════
content_slide("亚洲特殊性（1/2）", [
    "隐私合规：PIPL（中国）/ APPI（日本）/ PDPA（东南亚）按区域动态调整",
    "多语言支持：LanguageRouter 架构，原生支持中 / 日 / 韩 / 英四语",
    "合成 Persona：解决新市场冷启动问题，无需真实数据即可建模",
    "文化适配：不同市场的消费习惯、审美偏好、节日营销差异",
    "数据本地化：各区域数据存储满足当地法规要求",
], 15)

# ══════════════════════════════════════════════════════════════
# Page 16: 亚洲特殊性（2/2）
# ══════════════════════════════════════════════════════════════
content_slide("亚洲特殊性（2/2）", [
    "支付碎片化：MCP 统一支付网关，适配微信/支付宝/Line Pay 等",
    "城乡差异：StoreStrategy 分级策略，一线城市 vs 下沉市场",
    "电商生态：京东 / 天猫 / 乐天 / Shopee 多平台适配",
    "社交电商：小红书 / 抖音 / Instagram 内容种草链路",
    "物流差异：即时配送（中国）vs 标准配送（日本/东南亚）",
], 16)

# ══════════════════════════════════════════════════════════════
# Page 17: 业务场景 1 — 跨渠道个性化
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "业务场景 1 — 跨渠道个性化")

# Scenario
add_textbox(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.5),
            "场景：客户在 App 浏览了摇粒绒外套，到店后收到个性化推荐",
            font_size=18, color=DEEP_BLUE, bold=True)

# Flow steps
steps = [
    "① Customer Understanding Agent\n收集多渠道行为数据",
    "② Personalization Agent\n生成个性化推荐",
    "③ 门店 Agent\n推送搭配建议",
    "④ Marketing Agent\n发送优惠券",
]
for i, step in enumerate(steps):
    x = Inches(0.8 + i * 3.1)
    box = add_rounded_rect(slide, x, Inches(2.8), Inches(2.8), Inches(2.5), DEEP_BLUE)
    add_textbox(slide, x + Inches(0.15), Inches(3.0), Inches(2.5), Inches(2.2),
                step, font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)
    if i < 3:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
            x + Inches(2.8), Inches(3.7), Inches(0.3), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0),
            "预期效果：转化率提升 35%，客单价提升 20%，客户满意度提升至 85%",
            font_size=16, color=DARK_TEXT, bold=True)
add_page_number(slide, 17)

# ══════════════════════════════════════════════════════════════
# Page 18: 业务场景 2 — 智能库存管理
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "业务场景 2 — 智能库存管理")

add_textbox(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.5),
            "场景：季节交替时，自动预测需求并调整库存分配",
            font_size=18, color=DEEP_BLUE, bold=True)

steps = [
    "① Inventory Agent\n分析历史销售 + 天气数据",
    "② 需求预测\n7 天销量预测 + 安全库存",
    "③ 自动补货\n生成补货订单",
    "④ Pricing Agent\n联动调整促销价格",
]
for i, step in enumerate(steps):
    x = Inches(0.8 + i * 3.1)
    box = add_rounded_rect(slide, x, Inches(2.8), Inches(2.8), Inches(2.5), DEEP_BLUE)
    add_textbox(slide, x + Inches(0.15), Inches(3.0), Inches(2.5), Inches(2.2),
                step, font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)
    if i < 3:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
            x + Inches(2.8), Inches(3.7), Inches(0.3), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0),
            "预期效果：库存浪费减少 50%，缺货率降低 60%，周转率提升 40%",
            font_size=16, color=DARK_TEXT, bold=True)
add_page_number(slide, 18)

# ══════════════════════════════════════════════════════════════
# Page 19: 业务场景 3 — 动态定价
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "业务场景 3 — 动态定价")

add_textbox(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(0.5),
            "场景：竞品降价时，自动评估并调整价格策略",
            font_size=18, color=DEEP_BLUE, bold=True)

steps = [
    "① Pricing Agent\n监控竞品价格变动",
    "② 弹性分析\n价格弹性模型评估影响",
    "③ 策略生成\nThree-Layer Architecture决策",
    "④ 人工审批\n关键价格变动需确认",
]
for i, step in enumerate(steps):
    x = Inches(0.8 + i * 3.1)
    box = add_rounded_rect(slide, x, Inches(2.8), Inches(2.8), Inches(2.5), DEEP_BLUE)
    add_textbox(slide, x + Inches(0.15), Inches(3.0), Inches(2.5), Inches(2.2),
                step, font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)
    if i < 3:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
            x + Inches(2.8), Inches(3.7), Inches(0.3), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = GOLD
        arrow.line.fill.background()

add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0),
            "预期效果：利润率提升 15%，价格响应时间从 7 天缩短至 2 小时",
            font_size=16, color=DARK_TEXT, bold=True)
add_page_number(slide, 19)

# ══════════════════════════════════════════════════════════════
# Page 20: 预期效果
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "预期效果")

metrics = [
    ("+40%", "同店销售增长"),
    ("-50%", "库存浪费降低"),
    ("85%", "客户满意度"),
    ("+60%", "客户终身价值"),
    ("35%+", "电商渗透率"),
]
for i, (value, label) in enumerate(metrics):
    x = Inches(0.5 + i * 2.5)
    y = Inches(2.2)
    box = add_rounded_rect(slide, x, y, Inches(2.2), Inches(4.0), DEEP_BLUE)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.8), Inches(2.0), Inches(1.2),
                value, font_size=36, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), y + Inches(2.2), Inches(2.0), Inches(1.0),
                label, font_size=16, color=WHITE, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 20)

# ══════════════════════════════════════════════════════════════
# Page 21: 实施路径
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "实施路径")

phases = [
    ("Phase 1", "3 个月", "MVP 验证", "10 家门店试点\nCustomer Understanding + Personalization\n验证核心价值", DEEP_BLUE),
    ("Phase 2", "6 个月", "规模扩展", "100 家门店\n加入 Inventory + Pricing Agent\n全链路打通", MID_BLUE),
    ("Phase 3", "12 个月", "全渠道覆盖", "全渠道上线\n5 Agent 完整协同\n持续优化迭代", RGBColor(0x00, 0x55, 0x99)),
]
for i, (name, duration, title, desc, color) in enumerate(phases):
    x = Inches(0.8 + i * 4.2)
    box = add_rounded_rect(slide, x, Inches(2.0), Inches(3.8), Inches(4.8), color)
    add_textbox(slide, x + Inches(0.2), Inches(2.2), Inches(3.4), Inches(0.5),
                f"{name}（{duration}）", font_size=18, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), Inches(2.8), Inches(3.4), Inches(0.5),
                title, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_rect(slide, x + Inches(0.5), Inches(3.4), Inches(2.8), Inches(0.03), GOLD)
    add_textbox(slide, x + Inches(0.3), Inches(3.6), Inches(3.2), Inches(2.8),
                desc, font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 21)

# ══════════════════════════════════════════════════════════════
# Page 22: 风险与应对
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "风险与应对")

risks = [
    ("幻觉风险", "LLM 生成不准确信息", "RAG + 事实检查 + 人工审核"),
    ("延迟风险", "Agent 响应时间过长", "事件驱动 + 缓存 + 异步处理"),
    ("成本风险", "API 调用费用过高", "小模型处理简单任务 + 缓存策略"),
    ("合规风险", "数据隐私法规违规", "合成 Persona + 数据本地化 + 审计"),
]
for i, (risk, desc, solution) in enumerate(risks):
    y = Inches(1.7 + i * 1.35)
    # Risk box (red-ish)
    box = add_rounded_rect(slide, Inches(0.8), y, Inches(5.5), Inches(1.15), RGBColor(0xFF, 0xF0, 0xF0), RGBColor(0xCC, 0x33, 0x33))
    add_textbox(slide, Inches(1.0), y + Inches(0.1), Inches(5.0), Inches(0.4),
                f"⚠️ {risk}", font_size=16, color=RGBColor(0xCC, 0x33, 0x33), bold=True)
    add_textbox(slide, Inches(1.0), y + Inches(0.55), Inches(5.0), Inches(0.4),
                desc, font_size=13, color=DARK_TEXT)
    # Solution box (green-ish)
    box2 = add_rounded_rect(slide, Inches(6.8), y, Inches(5.8), Inches(1.15), RGBColor(0xF0, 0xFF, 0xF0), RGBColor(0x33, 0x99, 0x33))
    add_textbox(slide, Inches(7.0), y + Inches(0.1), Inches(5.4), Inches(0.4),
                "✅ 应对方案", font_size=16, color=RGBColor(0x33, 0x99, 0x33), bold=True)
    add_textbox(slide, Inches(7.0), y + Inches(0.55), Inches(5.4), Inches(0.4),
                solution, font_size=13, color=DARK_TEXT)
add_page_number(slide, 22)

# ══════════════════════════════════════════════════════════════
# Page 23: 成本估算
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "成本估算", "原型阶段月度成本")

# Table
headers = ["服务", "月度成本", "说明"]
rows = [
    ["Azure OpenAI", "~$500/月", "GPT-4o 推理调用"],
    ["Cosmos DB", "~$200/月", "全球分布式存储"],
    ["Azure ML", "~$300/月", "模型训练与部署"],
    ["其他服务", "~$200/月", "Event Hubs / APIM / Fabric"],
    ["总计", "~$1,300/月", "原型阶段预估"],
]
rows_count = len(rows) + 1
cols_count = len(headers)
table_shape = slide.shapes.add_table(rows_count, cols_count,
    Inches(2.5), Inches(2.0), Inches(8.3), Inches(4.0))
table = table_shape.table

# Header
for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = DEEP_BLUE
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

# Data rows
for r_idx, row in enumerate(rows, 1):
    for c_idx, val in enumerate(row):
        cell = table.cell(r_idx, c_idx)
        cell.text = val
        if r_idx == len(rows):  # Total row
            cell.fill.solid()
            cell.fill.fore_color.rgb = GOLD
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_TEXT if r_idx < len(rows) else DEEP_BLUE
            p.font.bold = (r_idx == len(rows))
            p.alignment = PP_ALIGN.CENTER

add_textbox(slide, Inches(2.5), Inches(6.2), Inches(8.3), Inches(0.5),
            "注：生产环境成本将随门店规模线性扩展，预计 ROI > 10x",
            font_size=13, color=RGBColor(0x66,0x66,0x66), alignment=PP_ALIGN.CENTER)
add_page_number(slide, 23)

# ══════════════════════════════════════════════════════════════
# Page 24: 总结
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, WHITE)
add_title_bar(slide, "总结")

summary_points = [
    "✓  5 Agent 协同解决零售业四大核心痛点",
    "✓  Azure 全栈技术支撑，企业级安全与可扩展性",
    "✓  亚洲市场特殊性深度适配，差异化竞争优势",
    "✓  可落地的三阶段实施路径，风险可控",
    "✓  预期 ROI > 10x，3 个月内验证核心价值",
]
add_bullet_textbox(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(3.5),
                   summary_points, font_size=20, color=DARK_TEXT, bullet_char='')

# Conclusion box
box = add_rounded_rect(slide, Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.5), DEEP_BLUE)
add_textbox(slide, Inches(1.3), Inches(5.4), Inches(10.7), Inches(1.1),
            "Project 46 不仅是一个技术方案，更是零售业数字化转型的范式革新。\n"
            "通过 AI Agent 协同，让优衣库从「经验驱动」迈向「数据智能驱动」。",
            font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 24)

# ══════════════════════════════════════════════════════════════
# Page 25: Thank You
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, DEEP_BLUE)

# Gold accent line at top
add_rect(slide, Inches(0), Inches(0), W, Inches(0.08), GOLD)

add_textbox(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5),
            "Thank You", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.5), Inches(11.3), Inches(0.8),
            "谢谢", font_size=32, color=GOLD, alignment=PP_ALIGN.CENTER)

# Divider
add_rect(slide, Inches(5), Inches(4.5), Inches(3.3), Inches(0.04), GOLD)

add_multiline_textbox(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(1.5), [
    "Q & A",
    "",
    "汇报人：圆子  |  Senior Solution Engineer",
    "Project 46 — Autonomous Retail Intelligence Agent Platform",
], font_size=18, color=RGBColor(0xCC,0xCC,0xCC), alignment=PP_ALIGN.CENTER)

# Bottom gold line
add_rect(slide, Inches(0), Inches(7.42), W, Inches(0.08), GOLD)

# ── Save ──────────────────────────────────────────────────────
output_path = "/home/azureadmin/.openclaw/workspace/ama-project/AMA-终汇报-PPT.pptx"
prs.save(output_path)
print(f"✅ PPT saved to {output_path}")
print(f"   Total slides: {len(prs.slides)}")
